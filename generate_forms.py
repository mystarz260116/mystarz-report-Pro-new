#!/usr/bin/env python3
"""
歯科技工所 日報入力用紙ジェネレーター
A5サイズ（縦）PowerPoint形式

使い方:
  python generate_forms.py
  → 日報入力用紙_A5.pptx が生成されます
"""

from pptx import Presentation
from pptx.util import Cm, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn
from lxml import etree

# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
# レイアウト定数
# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
A5_W = Cm(14.8)
A5_H = Cm(21.0)
M    = Cm(0.75)        # ページマージン（上下左右）
CW   = A5_W - 2 * M   # コンテンツ幅 ≈ 13.3 cm

# テーブル行高さ（EMU）
ROW_H     = int(Cm(0.52))   # 通常データ行
HDR_ROW_H = int(Cm(0.56))   # テーブルヘッダー行

# 備考欄
REMARKS_H     = Cm(1.75)
REMARKS_BOX_H = Cm(1.15)

# 1ページ目: ヘッダー帯(1.0) + 情報欄3行(3×0.68=2.04) + 余白(0.2) = 4.07cm
FIRST_TABLE_TOP = int(Cm(4.07))
FIRST_TABLE_BOT = int(A5_H - M - REMARKS_H)    # 備考欄を確保

# 続きページ
CONT_TABLE_TOP  = int(Cm(2.2))
CONT_TABLE_BOT  = int(A5_H - M - REMARKS_H)

# 各ページの最大データ行数
FIRST_MAX = (FIRST_TABLE_BOT - FIRST_TABLE_TOP - HDR_ROW_H) // ROW_H
CONT_MAX  = (CONT_TABLE_BOT  - CONT_TABLE_TOP  - HDR_ROW_H) // ROW_H


# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
# 部署データ（constants.ts より）
# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
DEPARTMENTS = [
    {
        'label': '大阪模型',
        'color': '3b82f6',
        'sections': [
            ('製作品目', [
                'ノーマル模型（急ぎ）', 'ノーマル模型【メタル】（総製作）', 'ノーマル模型【CAD】（総製作）',
                '貼り付け模型（急ぎ）', '貼り付け模型【メタル】（総製作）', '貼り付け模型【CAD】（総製作）',
                'インレー・コア模型（急ぎ）', 'インレー・コア模型（総製作）', '総数（急ぎ）', '総数（総製作）',
            ]),
        ],
    },
    {
        'label': 'メタル①',
        'color': '10b981',
        'sections': [
            ('WAX', ['HR', 'FMC', 'インレー', 'コア', 'その他']),
            ('埋没', ['スプルー植立', '埋没', 'キャスト', '割り出し', 'カット計量', 'パターン', '埋没 その他']),
        ],
    },
    {
        'label': 'メタル②',
        'color': '6366f1',
        'sections': [
            ('品目', ['クラウン', 'インレー', 'コア', '自費クラウン', '自費インレー', '自費コア']),
            ('調整・適合', [
                '調整・適合', '調整・コンタクト', '調整・バイト', 'ネジ付け・FMC/In',
                '研磨・FMC/In', '研磨・ブリッジ', '研磨・コア', 'ネジ外し・FMC/In', 'レーズ・ブリッジ',
            ]),
        ],
    },
    {
        'label': 'メタル③',
        'color': 'ec4899',
        'sections': [
            ('工程', [
                'メタル(適合～オペーク)', '築盛(築盛)', '形態(コンタクト～形態)',
                '研磨(シリコン～)', 'ホワイトWAX(築盛～形態)', 'トリミング(チェック～)',
            ]),
            ('製作品目', [
                'トリミング', 'ハードレジン', 'HJK', 'HB（インレー）', 'HB(アンレー)',
                'HB(ジャケット)', 'HB(金属裏装)', 'ファイバーコア(自費)', 'ファイバーコア(保険)',
                'CRインレー', 'CRアンレー', 'クラウン', 'インレー', 'ホワイトWAX',
            ]),
        ],
    },
    {
        'label': '自費',
        'color': '8b5cf6',
        'sections': [
            ('e.max', [
                'マウント(e.max)', 'トリミング(e.max)', 'パターン Cr(e.max)', 'パターン インレー(e.max)',
                'パターン ラミネート(e.max)', '植立：埋没・プレス(e.max)', '適合・調整(e.max)',
                'ステイン・完成 Cr(e.max)', 'ステイン・完成 インレー(e.max)', 'ステイン・完成 ラミネート(e.max)',
            ]),
            ('MB', [
                'マウント(MB)', 'トリミング(MB)', 'パターン(フルカントゥア)', 'パターン(キャップ)',
                '植立・埋没・キャスト(MB)', 'メタル調整(MB)', '前ロウ(ヶ所)', '築盛(MB)', '形態修正・完成(MB)',
            ]),
            ('Zirconia', [
                'マウント(Zir)', 'トリミング(Zir)', '設計・Cr(Zir)', '設計・In(Zir)', '適合・調整(Zir)',
                'フルジルコニア ステイン・完成(Cr)', 'フルジルコニア ステイン・完成(インレー)',
                'レイヤリング(築盛)(Zir)', 'レイヤリング(形成修正・完成)(Zir)',
            ]),
            ('インプラント', [
                '模型作り', 'マウント', 'アバットメント設計', 'アバットメント調整',
                'セメンテーション（スクリュー）', 'トランスファージグ',
            ]),
        ],
    },
    {
        'label': 'CAD/CAM①',
        'color': '22d3ee',
        'sections': [
            ('実績入力', ['トリミング', 'スキャン', 'CAM', '3Dプリンター']),
        ],
    },
    {
        'label': 'CAD/CAM②（設計）',
        'color': '06b6d4',
        'sections': [
            ('模型あり', ['模型あり 前歯', '模型あり クラウン', '模型あり インレー']),
            ('模型なし', ['模型なし 前歯', '模型なし クラウン', '模型なし インレー']),
            ('AI設計',   ['AI設計 前歯',   'AI設計 クラウン',   'AI設計 インレー']),
        ],
    },
    {
        'label': 'CAD/CAM③（完成）',
        'color': '0891b2',
        'sections': [
            ('実績入力', ['模型あり', '模型なし']),
        ],
    },
    {
        'label': 'デンチャー',
        'color': 'ef4444',
        'sections': [
            ('基本', ['台付(個)', 'トリミング(個)', 'バイト(ケース)', 'マウント(ケース)', '印象(個)']),
            ('3D/CAD', ['3Dデンチャー 設計(本)', '3Dデンチャー 完成(本)', 'CAD/CAM(設計)', 'CAD/CAM(完成)']),
            ('製作', [
                'ベース', 'ロー堤', 'トレー', 'クラスプ・バー設計(本)', 'クラスプ・バー パターン(本)',
                'クラスプ・バー 埋没(本)', 'クラスプ ワイヤー屈曲(本)', 'サンドブラスト(本)',
                'クラスプ・バー 研磨(本)', 'クラスプ適合(本)', 'ソルダーロー着(ヶ所)',
                'コバルトロー着(ヶ所)', '補強床※(枚)',
            ]),
            ('排列', [
                '試適排列(1-4歯)(床)', '試適排列(5-8歯)(床)', '試適排列(9-12歯)(床)',
                '試適排列(13-総義歯)(床)', '重合排列-咬合調整(1-4歯)(床)', '重合排列-咬合調整(5-8歯)(床)',
                '重合排列-咬合調整(9-12歯)(床)', '重合排列-咬合調整(13-総義歯)(床)',
                '組み立て(床)', '試適形成(床)', '重排形成(床)', '補強線屈曲(床)', '補強線ロー着(床)',
            ]),
            ('設計', ['デンチャー設計(ブロックアウト含む)(床)', '副模型製作(個)']),
            ('維持屈曲', [
                'デンチャー埋没(床)', '脱漏(自費は１人・保険２人当たりの数)(床)', '墳入(リング)',
                '割り出し　前工程(リング)', '対合はずし・洗浄(床)',
                '流し込み　前処理(シリコン型取り)(床)', '流し込み(人工歯置き換え含む)(ケース)',
            ]),
            ('プレス', [
                'カスタムトレー(ホワイトニング用)(床)', 'ナイトガードソフト(プレス)(床)',
                '咬合調整(ナイトガードソフトラミネート処理・有)(床)',
                '咬合調整(ナイトガードソフトラミネート処理・無)(床)',
                'ナイトガードハード(プレス)(床)', '咬合調整(ナイトガードハード)(床)',
                'スポーツマウスピース(ラミネート処理・有)(床)', 'スポーツマウスピース(ラミネート処理・無)(床)',
            ]),
            ('修理研磨', [
                'ソフトリライニング(墳入～仕上まで)(床)', '義歯修理(破折・増歯・補強線追加)(床)',
                '適合(床)', '床研磨(床)',
            ]),
            ('その他', [
                'ネーム入れ(デンチャー・プレス)(床)', 'メッシュプレート(動揺歯固定)(枚)',
                'バリオ(床)', 'チェック', '矯正', 'その他',
            ]),
        ],
    },
]


# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
# ユーティリティ関数
# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
def hex2rgb(h):
    h = h.lstrip('#')
    return int(h[:2], 16), int(h[2:4], 16), int(h[4:6], 16)


def lighten(c, pct):
    """色をpct(0.0-1.0)だけ白に近づける"""
    return tuple(min(255, int(v + (255 - v) * pct)) for v in c)


def set_cell_fill(cell, r, g, b):
    tc   = cell._tc
    tcPr = tc.get_or_add_tcPr()
    for tag in ('a:solidFill', 'a:noFill', 'a:gradFill', 'a:pattFill', 'a:blipFill'):
        for el in tcPr.findall(qn(tag)):
            tcPr.remove(el)
    sf = etree.SubElement(tcPr, qn('a:solidFill'))
    cl = etree.SubElement(sf, qn('a:srgbClr'))
    cl.set('val', f'{r:02X}{g:02X}{b:02X}')


def set_cell_text(cell, text, pt=8, bold=False, align=PP_ALIGN.LEFT, color=(25, 25, 25)):
    # セル内マージン
    tc   = cell._tc
    tcPr = tc.get_or_add_tcPr()
    tcPr.set('marT', '0')
    tcPr.set('marB', '0')
    tcPr.set('marL', str(int(Cm(0.08))))
    tcPr.set('marR', str(int(Cm(0.05))))

    # テキスト垂直中央
    tf     = cell.text_frame
    txBody = tf._txBody
    bodyPr = txBody.bodyPr
    bodyPr.set('anchor', 'ctr')

    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment    = align
    p.space_before = Pt(0)
    p.space_after  = Pt(0)

    for run in list(p.runs):
        p._p.remove(run._r)

    run = p.add_run()
    run.text           = text
    run.font.size      = Pt(pt)
    run.font.bold      = bold
    run.font.color.rgb = RGBColor(*color)


def set_table_borders(tbl, color_hex='C8C8C8'):
    for ri in range(len(tbl.rows)):
        for ci in range(len(tbl.columns)):
            tc   = tbl.cell(ri, ci)._tc
            tcPr = tc.get_or_add_tcPr()
            for side in ('lnL', 'lnR', 'lnT', 'lnB'):
                for old in tcPr.findall(qn(f'a:{side}')):
                    tcPr.remove(old)
                ln = etree.SubElement(tcPr, qn(f'a:{side}'))
                ln.set('w', '9525')
                ln.set('cap', 'flat')
                ln.set('cmpd', 'sng')
                sf = etree.SubElement(ln, qn('a:solidFill'))
                cl = etree.SubElement(sf, qn('a:srgbClr'))
                cl.set('val', color_hex)


def add_run(para, text, pt, bold=False, underline=False, color=(25, 25, 25)):
    run = para.add_run()
    run.text           = text
    run.font.size      = Pt(pt)
    run.font.bold      = bold
    run.font.underline = underline
    run.font.color.rgb = RGBColor(*color)


def make_textbox(slide, x, y, w, h):
    tb = slide.shapes.add_textbox(int(x), int(y), int(w), int(h))
    tf = tb.text_frame
    tf.word_wrap = False
    p  = tf.paragraphs[0]
    p.space_before = Pt(0)
    p.space_after  = Pt(0)
    return p


# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
# スライド構築
# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
def draw_header(slide, label, cr, cg, cb, page_idx, total_pages):
    hdr_h = Cm(1.0)
    s = slide.shapes.add_shape(1, int(M), int(M), int(CW), int(hdr_h))
    s.fill.solid()
    s.fill.fore_color.rgb = RGBColor(cr, cg, cb)
    s.line.fill.background()

    tf     = s.text_frame
    txBody = tf._txBody
    bodyPr = txBody.bodyPr
    bodyPr.set('anchor', 'ctr')

    tf.word_wrap = False
    p = tf.paragraphs[0]
    p.alignment    = PP_ALIGN.CENTER
    p.space_before = Pt(0)
    p.space_after  = Pt(0)

    suffix = f'  （{page_idx + 1}/{total_pages}）' if total_pages > 1 else ''
    add_run(p, f'【 {label} 】  日  報{suffix}', pt=13, bold=True, color=(255, 255, 255))


def draw_info(slide):
    """日付・時刻・担当者欄を追加。テーブル開始Y(int EMU)を返す。"""
    y0    = int(M + Cm(1.08))   # ヘッダー帯の下 + 余白
    row_h = int(Cm(0.68))
    fs    = 9.5

    # ─ 日付 ─
    p1 = make_textbox(slide, M, y0, CW, Cm(0.6))
    add_run(p1, '日付：', fs)
    add_run(p1, '          ', fs, underline=True)
    add_run(p1, ' 年 ', fs)
    add_run(p1, '       ', fs, underline=True)
    add_run(p1, ' 月 ', fs)
    add_run(p1, '       ', fs, underline=True)
    add_run(p1, ' 日', fs)

    # ─ 開始・終了時刻 ─
    p2 = make_textbox(slide, M, y0 + row_h, CW, Cm(0.6))
    add_run(p2, '開始時刻：', fs)
    add_run(p2, '       ', fs, underline=True)
    add_run(p2, ' ：', fs)
    add_run(p2, '       ', fs, underline=True)
    add_run(p2, '          終了時刻：', fs)
    add_run(p2, '       ', fs, underline=True)
    add_run(p2, ' ：', fs)
    add_run(p2, '       ', fs, underline=True)

    # ─ 担当者 ─
    p3 = make_textbox(slide, M, y0 + row_h * 2, CW, Cm(0.6))
    add_run(p3, '担当者：', fs)
    add_run(p3, '                                                             ', fs, underline=True)

    # 区切り線（薄い矩形）
    sep_y = y0 + row_h * 3 + int(Cm(0.05))
    sep   = slide.shapes.add_shape(1, int(M), sep_y, int(CW), int(Cm(0.03)))
    sep.fill.solid()
    sep.fill.fore_color.rgb = RGBColor(cr := 200, cg := 200, cb := 200)  # noqa
    sep.fill.fore_color.rgb = RGBColor(200, 200, 200)
    sep.line.fill.background()

    return y0 + row_h * 3 + int(Cm(0.15))   # テーブル開始Y


def draw_table(slide, rows, table_top, dept_rgb):
    cr, cg, cb = dept_rgb
    n = len(rows) + 1   # ヘッダー行 + データ行
    h = HDR_ROW_H + len(rows) * ROW_H

    col1_w = int(Cm(9.5))
    col2_w = int(CW) - col1_w

    tbl = slide.shapes.add_table(n, 2, int(M), int(table_top), int(CW), h).table
    tbl.columns[0].width = col1_w
    tbl.columns[1].width = col2_w

    # ── テーブルヘッダー行 ──
    tbl.rows[0].height = HDR_ROW_H
    set_cell_fill(tbl.cell(0, 0), cr, cg, cb)
    set_cell_fill(tbl.cell(0, 1), cr, cg, cb)
    set_cell_text(tbl.cell(0, 0), '項　　目', pt=8, bold=True, align=PP_ALIGN.CENTER, color=(255, 255, 255))
    set_cell_text(tbl.cell(0, 1), '数　量',   pt=8, bold=True, align=PP_ALIGN.CENTER, color=(255, 255, 255))

    sec_bg = lighten((cr, cg, cb), 0.80)   # セクションヘッダー行の背景色

    for i, (rtype, text) in enumerate(rows):
        ri = i + 1
        tbl.rows[ri].height = ROW_H
        if rtype == 'sec':
            set_cell_fill(tbl.cell(ri, 0), *sec_bg)
            set_cell_fill(tbl.cell(ri, 1), *sec_bg)
            set_cell_text(tbl.cell(ri, 0), f'▶  {text}', pt=7.5, bold=True)
            set_cell_text(tbl.cell(ri, 1), '',            pt=7.5)
        else:
            set_cell_fill(tbl.cell(ri, 0), 255, 255, 255)
            set_cell_fill(tbl.cell(ri, 1), 255, 255, 255)
            set_cell_text(tbl.cell(ri, 0), f'  {text}', pt=7.5)
            set_cell_text(tbl.cell(ri, 1), '',          pt=7.5)

    set_table_borders(tbl)


def draw_remarks(slide):
    rem_y = int(A5_H - M - REMARKS_H)

    p = make_textbox(slide, M, rem_y, Cm(2.0), Cm(0.45))
    add_run(p, '備　考：', 9)

    box_y = rem_y + int(Cm(0.45))
    box   = slide.shapes.add_shape(1, int(M), box_y, int(CW), int(REMARKS_BOX_H))
    box.fill.background()
    box.line.color.rgb = RGBColor(160, 160, 160)
    box.line.width     = Pt(0.75)


# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
# ページ生成
# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
def build_rows(dept):
    rows = []
    for title, items in dept['sections']:
        rows.append(('sec', title))
        for item in items:
            rows.append(('item', item))
    return rows


def add_dept_slides(prs, dept):
    cr, cg, cb = hex2rgb(dept['color'])
    all_rows   = build_rows(dept)

    # ページ分割
    pages     = []
    remaining = list(all_rows)
    while remaining:
        limit = FIRST_MAX if len(pages) == 0 else CONT_MAX
        pages.append(remaining[:limit])
        remaining = remaining[limit:]

    n_pages = len(pages)

    for pi, page_rows in enumerate(pages):
        is_first = (pi == 0)
        is_last  = (pi == n_pages - 1)

        slide = prs.slides.add_slide(prs.slide_layouts[6])   # Blankレイアウト

        draw_header(slide, dept['label'], cr, cg, cb, pi, n_pages)

        if is_first:
            table_top = draw_info(slide)
        else:
            p = make_textbox(slide, M, int(M + Cm(1.1)), CW, Cm(0.35))
            p.alignment = PP_ALIGN.RIGHT
            add_run(p, '（前ページからの続き）', pt=7.5, color=(120, 120, 120))
            table_top = CONT_TABLE_TOP

        draw_table(slide, page_rows, table_top, (cr, cg, cb))

        if is_last:
            draw_remarks(slide)


# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
# メイン
# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
def main():
    prs = Presentation()
    prs.slide_width  = A5_W
    prs.slide_height = A5_H

    print('Generating forms...\n')

    slide_info = []
    for dept in DEPARTMENTS:
        n_before = len(prs.slides)
        add_dept_slides(prs, dept)
        n_added  = len(prs.slides) - n_before
        slide_info.append((dept['label'], n_added))
        print(f'  OK {dept["label"]:20s}  {n_added}slide(s)')

    out = '日報入力用紙_A5.pptx'
    prs.save(out)
    total = sum(n for _, n in slide_info)
    print(f'\nDone: {out}  ({total} slides total)')
    print('Print setting: A5, fit to page, no margins')


if __name__ == '__main__':
    main()
